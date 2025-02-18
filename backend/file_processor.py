import os
import PyPDF2
import cv2
import pytesseract
from pptx import Presentation

def extract_text_from_file(uploaded_file):
    """
    Saves uploaded_file to disk and extracts text depending on file extension.
    """
    print(f"extract_text_from_file called for file: {uploaded_file.filename}") # ADDED
    os.makedirs("uploads", exist_ok=True)
    file_path = os.path.join("uploads", uploaded_file.filename)

    with open(file_path, "wb") as f:
        f.write(uploaded_file.file.read())

    ext = file_path.lower().split(".")[-1]

    if ext == "pdf":
        pdf_reader = PyPDF2.PdfReader(file_path)
        text = []
        for page in pdf_reader.pages:
            page_text = page.extract_text() or ""
            text.append(page_text)
        extracted_text = "\n".join(text).strip()
        print(f"Extracted text (pdf): {extracted_text}") # ADDED
        return extracted_text

    elif ext == "pptx":
        prs = Presentation(file_path)
        all_text = []
        for slide in prs.slides:
            # Extract from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text.append(shape.text)
            # Extract from notes if desired
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                all_text.append(slide.notes_slide.notes_text_frame.text)
        extracted_text = "\n".join(all_text).strip()
        print(f"Extracted text (pptx): {extracted_text}") # ADDED
        return extracted_text

    elif ext in ["png", "jpg", "jpeg"]:
        img = cv2.imread(file_path)
        text = pytesseract.image_to_string(img)
        extracted_text = text.strip()
        print(f"Extracted text (image): {extracted_text}") # ADDED
        return extracted_text

    else:
        raise ValueError("Unsupported file type")